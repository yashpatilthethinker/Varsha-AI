import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import io
import re
from typing import Dict, Any

class DocumentProcessor:
    """Handles document processing for various file formats"""
    
    def __init__(self):
        self.supported_formats = {
            'application/pdf': self._process_pdf,
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': self._process_docx,
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': self._process_pptx,
            'text/plain': self._process_txt
        }
    
    def process_file(self, uploaded_file) -> str:
        """Process uploaded file based on its type"""
        file_type = uploaded_file.type
        
        if file_type not in self.supported_formats:
            raise ValueError(f"Unsupported file type: {file_type}")
        
        # Read file content
        file_content = uploaded_file.read()
        
        # Reset file pointer for potential re-reading
        uploaded_file.seek(0)
        
        # Process based on file type
        processor = self.supported_formats[file_type]
        extracted_text = processor(file_content)
        
        # Clean and return text
        return self._clean_text(extracted_text)
    
    def _process_pdf(self, file_content: bytes) -> str:
        """Extract text from PDF files"""
        try:
            # Create PDF document from bytes
            pdf_document = fitz.open(stream=file_content, filetype="pdf")
            text_content = []
            
            for page_num in range(pdf_document.page_count):
                page = pdf_document.load_page(page_num)
                text = page.get_text()
                if text.strip():  # Only add non-empty pages
                    text_content.append(f"--- Page {page_num + 1} ---\n{text}")
            
            pdf_document.close()
            return "\n\n".join(text_content)
        
        except Exception as e:
            raise Exception(f"Error processing PDF: {str(e)}")
    
    def _process_docx(self, file_content: bytes) -> str:
        """Extract text from Word documents"""
        try:
            # Create document from bytes
            doc = Document(io.BytesIO(file_content))
            text_content = []
            
            # Extract paragraph text
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            # Extract table text
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        table_text.append(" | ".join(row_text))
                
                if table_text:
                    text_content.append("--- Table ---\n" + "\n".join(table_text))
            
            return "\n\n".join(text_content)
        
        except Exception as e:
            raise Exception(f"Error processing Word document: {str(e)}")
    
    def _process_pptx(self, file_content: bytes) -> str:
        """Extract text from PowerPoint presentations"""
        try:
            # Create presentation from bytes
            prs = Presentation(io.BytesIO(file_content))
            text_content = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"--- Slide {slide_num} ---"]
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                # Extract notes
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_text.append(f"Notes: {notes_text}")
                
                if len(slide_text) > 1:  # Only add if there's content beyond the header
                    text_content.append("\n".join(slide_text))
            
            return "\n\n".join(text_content)
        
        except Exception as e:
            raise Exception(f"Error processing PowerPoint: {str(e)}")
    
    def _process_txt(self, file_content: bytes) -> str:
        """Extract text from plain text files"""
        try:
            # Decode bytes to string
            text = file_content.decode('utf-8')
            return text
        
        except UnicodeDecodeError:
            # Try other encodings
            try:
                text = file_content.decode('latin-1')
                return text
            except Exception as e:
                raise Exception(f"Error decoding text file: {str(e)}")
        
        except Exception as e:
            raise Exception(f"Error processing text file: {str(e)}")
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize extracted text"""
        if not text:
            return ""
        
        # Remove excessive whitespace
        text = re.sub(r'\n\s*\n\s*\n', '\n\n', text)  # Reduce multiple newlines
        text = re.sub(r' {2,}', ' ', text)  # Reduce multiple spaces
        
        # Remove special characters that might cause issues
        text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', text)
        
        return text.strip()
    
    def get_document_info(self, uploaded_file) -> Dict[str, Any]:
        """Get basic information about the document"""
        return {
            'name': uploaded_file.name,
            'type': uploaded_file.type,
            'size': uploaded_file.size,
            'supported': uploaded_file.type in self.supported_formats
        }
    
    def extract_metadata(self, uploaded_file, file_content: bytes) -> Dict[str, Any]:
        """Extract metadata from documents where possible"""
        metadata = {
            'filename': uploaded_file.name,
            'file_type': uploaded_file.type,
            'file_size': uploaded_file.size
        }
        
        if uploaded_file.type == 'application/pdf':
            try:
                pdf_document = fitz.open(stream=file_content, filetype="pdf")
                pdf_metadata = pdf_document.metadata
                metadata.update({
                    'title': pdf_metadata.get('title', ''),
                    'author': pdf_metadata.get('author', ''),
                    'subject': pdf_metadata.get('subject', ''),
                    'page_count': pdf_document.page_count
                })
                pdf_document.close()
            except Exception:
                pass  # Metadata extraction failed, but that's okay
        
        return metadata